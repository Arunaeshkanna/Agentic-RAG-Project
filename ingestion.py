import os
from pptx import Presentation
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import (
    PyPDFLoader,
    Docx2txtLoader,
    TextLoader,
    WebBaseLoader
)

def load_documents(files, urls):
    documents = []

    for file in files:
        ext = file.name.split(".")[-1].lower()
        folder = f"data/{ext}"
        os.makedirs(folder, exist_ok=True)

        path = f"{folder}/{file.name}"
        with open(path, "wb") as f:
            f.write(file.read())

        if ext == "pdf":
            documents.extend(PyPDFLoader(path).load())
        elif ext in ["doc", "docx"]:
            documents.extend(Docx2txtLoader(path).load())
        elif ext == "txt":
            documents.extend(TextLoader(path).load())
        elif ext == "pptx":
            prs = Presentation(path)
            text = "\n".join(
                shape.text
                for slide in prs.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )
            if text.strip():
                documents.append(
                    Document(page_content=text,
                             metadata={"source": file.name})
                )

    for url in urls:
        url = url.strip()
        if url:
            documents.extend(WebBaseLoader(url).load())

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=300,
        chunk_overlap=100
    )

    return splitter.split_documents(documents)